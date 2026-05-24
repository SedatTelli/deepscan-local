"""
DeepScan Local - File Content Parser
Designer: Sedat Telli | sedattelli.com

Extracts raw text and its Turkish-normalized form from supported file types.
Uses Python 3.10+ match-case for clean extension dispatch.
"""

import re
from pathlib import Path

from config import MAX_FILE_BYTES, METADATA_ONLY_EXTENSIONS, log_error

# ---------------------------------------------------------------------------
# Turkish character normalization map
# Used to build content_normalized for FTS5 indexing.
# ---------------------------------------------------------------------------
_TR_MAP = str.maketrans(
    "ışğüöçİŞĞÜÖÇ",
    "isguocISGUOC"
)


def normalize(text: str) -> str:
    """Lowercase + apply Turkish → ASCII character substitution."""
    return text.lower().translate(_TR_MAP)


# ---------------------------------------------------------------------------
# Public extraction entry point
# ---------------------------------------------------------------------------

def extract(file_path: Path) -> tuple[str, str]:
    """
    Return (content_raw, content_normalized) for *file_path*.

    content_raw        — original text for snippet display
    content_normalized — lowercased + Turkish-mapped for FTS5 search

    Returns ("", "") on any failure or unsupported type.
    """
    suffix = file_path.suffix.lower()

    try:
        stat = file_path.stat()
        if stat.st_size == 0:
            return "", ""
        # Bypass the size cap for extensions where we never read file bytes
        if stat.st_size > MAX_FILE_BYTES and suffix not in METADATA_ONLY_EXTENSIONS:
            return "", ""
    except OSError:
        return "", ""

    raw = ""

    try:
        match suffix:
            case ".txt" | ".md" | ".csv" | ".log":
                raw = _read_text(file_path)

            case ".json" | ".xml" | ".svg":
                raw = _read_text(file_path)

            case ".py" | ".js":
                raw = _read_text(file_path)

            case ".pdf":
                raw = _extract_pdf(file_path)

            case ".docx":
                raw = _extract_docx(file_path)

            case ".xlsx":
                raw = _extract_xlsx(file_path)

            case ".pptx":
                raw = _extract_pptx(file_path)

            case ".psd" | ".ai":
                raw = _extract_adobe_xmp(file_path)

            case ".jpg" | ".jpeg" | ".png" | ".gif" | ".bmp" | ".webp" | ".tiff" | ".tif":
                raw = _extract_image_meta(file_path)
                # Try OCR if no EXIF text was found
                if not raw.strip() or len(raw.strip()) < 10:
                    ocr_text = _extract_image_ocr(file_path)
                    if ocr_text:
                        raw = ocr_text

            case ".mp4" | ".avi" | ".mkv" | ".mov" | ".wmv" | ".flv" | ".webm" | ".m4v":
                raw = _extract_video_meta(file_path)

            case _:
                return "", ""

    except Exception as exc:
        log_error(f"parser.extract [{suffix}] {file_path}: {exc}")
        return "", ""

    # Hard cap: keep first 400 000 chars to avoid huge DB rows
    raw = raw[:400_000].strip()
    if not raw:
        return "", ""

    return raw, normalize(raw)


# ---------------------------------------------------------------------------
# Snippet helper (for UI display)
# ---------------------------------------------------------------------------

def get_snippet(content: str, query: str, max_lines: int = 5) -> str:
    """
    Return the first paragraph that contains *query*, capped to *max_lines* lines.
    Falls back to the beginning of the content if no match is found.
    """
    if not content:
        return ""

    # Build ordered search terms: full phrase first, then individual words
    query_lower = query.lower().strip()
    terms = [query_lower] + [w for w in query_lower.split() if len(w) >= 2]

    if not terms or not terms[0]:
        lines = [l for l in content.splitlines() if l.strip()]
        return "\n".join(lines[:max_lines])

    # Split into paragraphs (double-newline separated)
    paragraphs = [p.strip() for p in content.split("\n\n") if p.strip()]
    # Excel/Word/PDF use single \n — no paragraph breaks, so fall back to line-by-line
    if len(paragraphs) <= 1:
        paragraphs = [l.strip() for l in content.splitlines() if l.strip()]

    # Find the first paragraph that contains any search term
    for term in terms:
        for para in paragraphs:
            if term in para.lower():
                lines = para.splitlines()
                capped = "\n".join(lines[:max_lines])
                # Hard char cap so wrapping doesn't overflow beyond 5 visual lines
                if len(capped) > 420:
                    capped = capped[:420].rstrip() + "…"
                return capped

    # Fallback: beginning of content (non-empty lines only)
    non_empty = [l.strip() for l in content.splitlines() if l.strip()]
    return "\n".join(non_empty[:max_lines])


# ---------------------------------------------------------------------------
# Format-specific extractors
# ---------------------------------------------------------------------------

def _read_text(path: Path) -> str:
    """Read plain-text file with multi-encoding fallback."""
    for enc in ("utf-8", "utf-8-sig", "cp1254", "cp1252", "latin-1"):
        try:
            return path.read_text(encoding=enc)
        except (UnicodeDecodeError, OSError):
            continue
    return ""


def _extract_pdf(path: Path) -> str:
    try:
        from pypdf import PdfReader
        reader = PdfReader(str(path))
        pages = []
        for page in reader.pages:
            text = page.extract_text()
            if text:
                pages.append(text)
        return "\n".join(pages)
    except Exception as exc:
        log_error(f"PDF {path}: {exc}")
        return ""


def _extract_docx(path: Path) -> str:
    try:
        from docx import Document
        doc = Document(str(path))
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        # Also extract from tables
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    if cell.text.strip():
                        paragraphs.append(cell.text)
        return "\n".join(paragraphs)
    except Exception as exc:
        log_error(f"DOCX {path}: {exc}")
        return ""


def _extract_xlsx(path: Path) -> str:
    try:
        import openpyxl
        wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
        rows: list[str] = []
        for sheet in wb.worksheets:
            for row in sheet.iter_rows(values_only=True):
                line = " ".join(str(c) for c in row if c is not None)
                if line.strip():
                    rows.append(line)
        wb.close()
        return "\n".join(rows)
    except Exception as exc:
        log_error(f"XLSX {path}: {exc}")
        return ""


def _extract_pptx(path: Path) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(str(path))
        parts: list[str] = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    parts.append(shape.text)
        return "\n".join(parts)
    except Exception as exc:
        log_error(f"PPTX {path}: {exc}")
        return ""


def _extract_image_meta(path: Path) -> str:
    """
    Try to pull EXIF text fields from JPEG images via Pillow (already a dep).
    Falls back to filename stem for all other image formats or on any error.
    """
    stem = path.stem
    if path.suffix.lower() not in (".jpg", ".jpeg"):
        return stem
    try:
        from PIL import Image as _PILImage
        from PIL.ExifTags import TAGS as _TAGS
        _TEXT_TAGS = frozenset({
            "Make", "Model", "Artist", "Copyright",
            "ImageDescription", "Software", "DateTime", "DateTimeOriginal",
        })
        with _PILImage.open(str(path)) as img:
            raw_exif = getattr(img, "_getexif", lambda: None)()
            if not raw_exif:
                return stem
            parts = [stem]
            for tag_id, val in raw_exif.items():
                tag = _TAGS.get(tag_id, "")
                if tag in _TEXT_TAGS and isinstance(val, str) and val.strip():
                    parts.append(val.strip())
            return " ".join(parts)
    except Exception:
        return stem


def _extract_video_meta(path: Path) -> str:
    """Index video files by their filename only (no file bytes are read)."""
    return path.stem


def _extract_image_ocr(path: Path) -> str:
    """
    Attempt OCR on an image using Windows.Media.Ocr (built-in, no extra deps)
    with a pytesseract fallback.  Returns "" gracefully if neither is available.
    """
    # ── Option 1: Windows built-in OCR (Windows 10 1809+) via winrt ─────────
    try:
        import asyncio
        import winrt.windows.media.ocr as _ocr
        import winrt.windows.storage as _storage
        import winrt.windows.storage.streams as _streams
        import winrt.windows.graphics.imaging as _imaging

        async def _run_ocr() -> str:
            engine = _ocr.OcrEngine.try_create_from_user_profile_languages()
            if engine is None:
                return ""
            file = await _storage.StorageFile.get_file_from_path_async(str(path.resolve()))
            stream = await file.open_async(_storage.FileAccessMode.READ)
            decoder = await _imaging.BitmapDecoder.create_async(stream)
            bmp = await decoder.get_software_bitmap_async()
            result = await engine.recognize_async(bmp)
            return result.text

        text = asyncio.run(_run_ocr())
        if text and text.strip():
            return text.strip()
    except Exception:
        pass

    # ── Option 2: pytesseract (requires tesseract-ocr installed separately) ──
    try:
        import pytesseract
        from PIL import Image as _PILImage
        with _PILImage.open(str(path)) as img:
            text = pytesseract.image_to_string(img, lang="tur+eng")
        if text and text.strip():
            return text.strip()
    except Exception:
        pass

    return ""


def _extract_adobe_xmp(path: Path) -> str:
    """
    Extract XMP metadata from PSD / AI files via regex on the first 64 KB.
    XMP is stored as UTF-8 XML between <x:xmpmeta> tags in the binary header.
    No third-party Adobe library required.
    """
    try:
        with path.open("rb") as fh:
            header = fh.read(65_536)

        match = re.search(
            b"<x:xmpmeta[^>]*>.*?</x:xmpmeta>",
            header,
            re.DOTALL,
        )
        if not match:
            return ""

        xmp_xml = match.group(0).decode("utf-8", errors="ignore")
        # Strip XML tags; keep human-readable text values
        text = re.sub(r"<[^>]+>", " ", xmp_xml)
        return " ".join(text.split())

    except Exception as exc:
        log_error(f"Adobe XMP {path}: {exc}")
        return ""
